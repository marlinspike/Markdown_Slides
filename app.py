from pptx.util import Pt, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
import os
import re
import argparse
from dotenv import load_dotenv

# Load environment settings for font sizes and bold setting
load_dotenv()
font_sizes = {
    'level-1': os.getenv('level-one', '18'),
    'level-2': os.getenv('level-two', '16'),
    'level-3': os.getenv('level-three', '14')
}
is_bold_level_one = os.getenv('bold-level-one', 'true').lower() == 'true'  # Default to bold for level-1

def get_font_size(level):
    """
    Fetches the font size from the environment settings based on the indentation level.
    """
    key = f'level-{level+1}'
    font_size = font_sizes.get(key, '12')
    return Pt(int(font_size))


def parse_markdown(md_file):
    """
    Parses markdown file to slides and bullet points, considering indentation levels and additional Markdown syntax.
    Uses 2 spaces per indentation level.
    """
    slides = []
    with open(md_file, 'r') as file:
        md_text = file.read()
    current_slide_number = 0
    current_slide = None
    for line in md_text.strip().split('\n'):
        slide_match = re.match(r'^Slide\s*(\d*)\s*:\s*(.*)', line)
        if slide_match:
            slide_number = slide_match.group(1)
            slide_title = slide_match.group(2).strip()
            current_slide_number += 1
            if slide_number and int(slide_number) == current_slide_number:
                current_slide_number = int(slide_number)
            current_slide = {'number': current_slide_number, 'title': slide_title, 'content': []}
            slides.append(current_slide)
        elif current_slide:
            leading_spaces = len(line) - len(line.lstrip(' '))
            level = leading_spaces // 2
            text = line.strip()
            if text.startswith('- '):
                text = text[2:]
            # Parse headings
            heading_match = re.match(r'^(#+) (.+)', text)
            if heading_match:
                heading_level = len(heading_match.group(1))
                text = heading_match.group(2)
                level = heading_level - 1
            # Parse formatting tags
            text = re.sub(r'<b>(.+?)</b>', r'**\1**', text)
            text = re.sub(r'<i>(.+?)</i>', r'_\1_', text)
            text = re.sub(r'<u>(.+?)</u>', r'__\1__', text)
            text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
            text = re.sub(r'_(.+?)_', r'<i>\1</i>', text)
            text = re.sub(r'__(.+?)__', r'<u>\1</u>', text)
            text = re.sub(r'~~(.+?)~~', r'<s>\1</s>', text)
            text = re.sub(r'`(.+?)`', r'<code>\1</code>', text)
            # Parse hyperlinks
            text = re.sub(r'\[(.+?)\]\((.+?)\)', r'<a href="\2">\1</a>', text)
            # Parse images
            text = re.sub(r'!\[(.+?)\]\((.+?)\)', r'<img src="\2" alt="\1">', text)
            current_slide['content'].append((text, level))
    return slides



def create_presentation(slides, pptx_file):
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]  # 'Title and Content' layout

    for slide_info in slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = slide_info['title']
        content_placeholder = slide.placeholders[1]

        for text, level in slide_info['content']:
            if level == 0:
                # Render heading level 1
                p = content_placeholder.text_frame.add_paragraph()
                p.text = text
                p.font.size = get_font_size(level)
                p.font.bold = True
            else:
                # Render bullet points for other heading levels and non-heading text
                p = content_placeholder.text_frame.add_paragraph()
                p.text = text
                p.level = level - 1  # Adjust the level for bullet points
                p.font.size = get_font_size(level)

            # Parse and apply formatting
            if '<b>' in text:
                for run in p.runs:
                    if '<b>' in run.text:
                        run.text = run.text.replace('<b>', '').replace('</b>', '')
                        run.font.bold = True
            if '<i>' in text:
                for run in p.runs:
                    if '<i>' in run.text:
                        run.text = run.text.replace('<i>', '').replace('</i>', '')
                        run.font.italic = True
            if '<u>' in text:
                for run in p.runs:
                    if '<u>' in run.text:
                        run.text = run.text.replace('<u>', '').replace('</u>', '')
                        run.font.underline = True
            if '<s>' in text:
                for run in p.runs:
                    if '<s>' in run.text:
                        run.text = run.text.replace('<s>', '').replace('</s>', '')
                        run.font.strike = True
            if '<code>' in text:
                for run in p.runs:
                    if '<code>' in run.text:
                        run.text = run.text.replace('<code>', '').replace('</code>', '')
                        run.font.name = 'Courier New'
            if '<a href=' in text:
                for run in p.runs:
                    if '<a href=' in run.text:
                        url = re.search(r'<a href="(.+?)">', run.text).group(1)
                        run.text = re.sub(r'<a href=".+?">', '', run.text).replace('</a>', '')
                        run.hyperlink.address = url
                        run.font.color.rgb = RGBColor(0, 0, 255)  # Blue color for hyperlinks
            if '<img src=' in text:
                for run in p.runs:
                    if '<img src=' in run.text:
                        src = re.search(r'<img src="(.+?)"', run.text).group(1)
                        alt = re.search(r'alt="(.+?)"', run.text).group(1)
                        run.text = ''
                        pic = slide.shapes.add_picture(src, Inches(1), Inches(1))
                        pic.left = Inches(1)
                        pic.top = Inches(1)

    prs.save(pptx_file)
    print(f"Presentation saved as {pptx_file}")


def process_markdown_file(md_file, output_dir):
    slides = parse_markdown(md_file)
    pptx_file = os.path.join(output_dir, os.path.splitext(os.path.basename(md_file))[0] + '.pptx')
    create_presentation(slides, pptx_file)


def main():
    parser = argparse.ArgumentParser(description="Convert Markdown files to PowerPoint presentations")
    parser.add_argument('-r', '--recursive', action='store_true', help='Process all Markdown files in the input folder recursively')
    parser.add_argument('-f', '--file', type=str, help='Markdown file to process (naked filename, relative path, or full path)')
    parser.add_argument('-o', '--output', type=str, help='Output PPTX filename (only applicable with -f option)')
    args = parser.parse_args()

    input_dir = 'input'
    output_dir = 'output'

    # Create the output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    if args.recursive:
        # Process all Markdown files in the input folder recursively
        for root, dirs, files in os.walk(input_dir):
            for file in files:
                if file.endswith('.md'):
                    md_file = os.path.join(root, file)
                    process_markdown_file(md_file, output_dir)
    elif args.file:
        # Process a single Markdown file specified by -f option
        md_file = args.file
        if not os.path.isabs(md_file):
            md_file = os.path.join(input_dir, md_file)
        if args.output:
            pptx_file = args.output
        else:
            pptx_file = os.path.join(output_dir, os.path.splitext(os.path.basename(md_file))[0] + '.pptx')
        process_markdown_file(md_file, os.path.dirname(pptx_file))
    else:
        print("Please specify either -r or -f option.")

if __name__ == '__main__':
    main()